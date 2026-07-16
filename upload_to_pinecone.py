import os, glob, hashlib, warnings, base64, time
import fitz  # PyMuPDF
import docx2txt
from pptx import Presentation
import pandas as pd
warnings.filterwarnings("ignore")

from pinecone import Pinecone
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openai import OpenAI

# 🔑 ตั้งค่า API KEY
PINECONE_API_KEY = "xxx"
OPENAI_API_KEY = "xxx"
INDEX_NAME = "atc-master"
PDF_FOLDER_PATH = r"D:\01 Project\Atc_Essential"

# 🌟 สวิตช์เปิด-ปิด โหมดอ่านภาพสแกน (OCR)
ENABLE_AI_OCR = False

client = OpenAI(api_key=OPENAI_API_KEY)

def get_hash_id(filename, chunk_index):
    return hashlib.md5(f"{filename}_{chunk_index}".encode()).hexdigest()

def extract_text_from_file(file_path):
    ext = file_path.lower().split('.')[-1]
    text = ""
    try:
        if ext == 'pdf':
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text().strip()
                
                # 🌟 ถ้าระบบเปิด OCR และหน้ากระดาษไม่มีตัวหนังสือ (เป็นภาพสแกน)
                if ENABLE_AI_OCR and len(page_text) < 50:
                    print(f"   🔍 กำลังทำ OCR หน้า {page_num + 1}/{len(doc)}...")
                    pix = page.get_pixmap(dpi=150)
                    img_base64 = base64.b64encode(pix.tobytes("jpeg")).decode('utf-8')
                    
                    # 🌟 ระบบวนลูปพยายามใหม่ (Retry) หากโดนบล็อก Rate Limit
                    success = False
                    while not success:
                        try:
                            ocr_res = client.chat.completions.create(
                                model="gpt-4o-mini", 
                                messages=[{"role": "user", "content": [
                                    {"type": "text", "text": "Extract all readable text from this document image in its original language. Do not add any extra comments."}, 
                                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_base64}"}}
                                ]}],
                                max_tokens=1500
                            )
                            page_text = ocr_res.choices[0].message.content
                            success = True # สำเร็จแล้ว ออกจากลูป
                            
                            # พัก 2 วินาที ป้องกัน Token เต็มเร็วไป
                            time.sleep(2) 
                            
                        except Exception as api_err:
                            if "429" in str(api_err):
                                print("   ⚠️ Rate Limit เต็ม! กำลังหยุดพัก 20 วินาที แล้วจะเริ่มทำหน้าเดิมใหม่...")
                                time.sleep(20)
                            else:
                                print(f"   ❌ ข้ามหน้า {page_num+1} เนื่องจาก Error อื่น: {api_err}")
                                break # ถ้าเป็น Error อื่นให้ข้ามไปเลย
                    
                text += f"\n{page_text}\n"
                
        elif ext in ['doc', 'docx']: text = docx2txt.process(file_path)
        elif ext in ['ppt', 'pptx']: 
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        elif ext in ['xls', 'xlsx']:
            df_dict = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, df in df_dict.items(): text += f"\n--- Sheet: {sheet_name} ---\n{df.to_string()}"
    except Exception as e: print(f"⚠️ อ่านไฟล์ไม่ได้: {file_path} - {e}")
    return text

def upload_to_pinecone():
    all_files = []
    for pattern in ['*.pdf', '*.doc', '*.docx', '*.ppt', '*.pptx', '*.xls', '*.xlsx']:
        all_files.extend(glob.glob(os.path.join(PDF_FOLDER_PATH, '**', pattern), recursive=True))
    
    if not all_files: return print("❌ ไม่พบไฟล์เอกสาร")
    
    pc = Pinecone(api_key=PINECONE_API_KEY)
    
    if INDEX_NAME not in pc.list_indexes().names():
        print("🏗️ สร้างฐานข้อมูลใหม่ ขนาด 768 มิติ...")
        from pinecone import ServerlessSpec
        pc.create_index(name=INDEX_NAME, dimension=768, metric="cosine", spec=ServerlessSpec(cloud="aws", region="us-east-1"))
        time.sleep(10)
    
    index = pc.Index(INDEX_NAME)
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)

    for i, file_path in enumerate(all_files, 1):
        try:
            filename = os.path.basename(file_path)
            folder_name = os.path.basename(os.path.dirname(file_path))
            namespace = "General" if folder_name == os.path.basename(PDF_FOLDER_PATH) else folder_name
                
            print(f"🚀 [{i}/{len(all_files)}] กำลังอัปโหลด: {filename}")
            raw_text = extract_text_from_file(file_path)
            if not raw_text.strip(): 
                print("   ⏩ ข้ามไฟล์นี้ (ไม่พบข้อความ)")
                continue
            
            splits = text_splitter.split_text(raw_text)
            vectors_to_upsert = []
            
            batch_size = 100
            for b_idx in range(0, len(splits), batch_size):
                batch_chunks = splits[b_idx : b_idx + batch_size]
                
                # 🌟 ระบบ Retry ตอนส่งเข้า Pinecone
                success_emb = False
                while not success_emb:
                    try:
                        res = client.embeddings.create(input=batch_chunks, model="text-embedding-3-small", dimensions=768)
                        for j, emb in enumerate(res.data):
                            vector_id = get_hash_id(filename, b_idx + j)
                            vectors_to_upsert.append({"id": vector_id, "values": emb.embedding, "metadata": {"text": batch_chunks[j], "source": filename}})
                        success_emb = True
                        time.sleep(1) # พักนิดนึง
                    except Exception as emb_err:
                        if "429" in str(emb_err):
                            print("   ⚠️ Embedding Rate Limit เต็ม! พัก 10 วินาที...")
                            time.sleep(10)
                        else:
                            break
            
            for j in range(0, len(vectors_to_upsert), 100):
                index.upsert(vectors=vectors_to_upsert[j:j+100], namespace=namespace)
            print("   ✅ อัปโหลดสำเร็จ!")
            
        except Exception as e: print(f"⚠️ Error ไฟล์นี้: {e}")
        
    print("\n🎉 อัปโหลดสำเร็จ 100%!")

if __name__ == "__main__":
    upload_to_pinecone()